from __future__ import annotations

from pathlib import Path
from subprocess import run

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
ASSETS = DOCS / "assets"
POSTHOC = ROOT / "reports" / "aggregate" / "posthoc_full_with_ppo_fix_20260326"
OUT = DOCS / "group_meeting_summary_20260326.pptx"
ARCH_SVG = ASSETS / "windblown_framework_architecture.svg"
ARCH_PNG = ASSETS / "windblown_framework_architecture.png"


def ensure_architecture_png() -> Path:
    if ARCH_PNG.exists() and ARCH_PNG.stat().st_mtime >= ARCH_SVG.stat().st_mtime:
        return ARCH_PNG
    run([
        "convert",
        "-background",
        "white",
        str(ARCH_SVG),
        str(ARCH_PNG),
    ], check=True)
    return ARCH_PNG


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(15, 23, 42)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = RGBColor(71, 85, 105)


def add_bullet_box(slide, x, y, w, h, title: str, bullets: list[str], fill=(248, 250, 252), line=(203, 213, 225)) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*line)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    r0.font.bold = True
    r0.font.size = Pt(15)
    r0.font.color.rgb = RGBColor(15, 23, 42)
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.bullet = True
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(30, 41, 59)


def add_caption(slide, x, y, w, text: str) -> None:
    box = slide.shapes.add_textbox(x, y, w, Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9.5)
    r.font.color.rgb = RGBColor(71, 85, 105)


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    arch_png = ensure_architecture_png()
    summary = pd.read_csv(ROOT / "reports" / "aggregate" / "metrics_forecast_all_full_with_ppo_fix_20260326_scheduler_summary.csv")
    primary = pd.read_csv(POSTHOC / "task_focus_primary" / "primary_target_scheduler_summary.csv")

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "功率约束下的多传感器调度与微气候预测",
        "问题定义、系统结构与实验主线",
    )
    slide.shapes.add_picture(str(arch_png), Inches(0.55), Inches(1.05), width=Inches(7.25))
    add_bullet_box(
        slide,
        Inches(8.0), Inches(1.1), Inches(4.7), Inches(2.0),
        "这套框架在做什么",
        [
            "同一条 truth 序列上比较不同调度策略，保证实验公平。",
            "调度结果先经过 OnlineSubsetProjector，再满足瞬时功率与启动峰值约束。",
            "观测先经 Kalman estimator，再生成 scheduler-specific forecasting dataset。",
            "最终比较的是：在节电前提下，下游预测还能保住多少性能。",
        ],
        fill=(239, 246, 255), line=(147, 197, 253),
    )
    add_bullet_box(
        slide,
        Inches(8.0), Inches(3.35), Inches(4.7), Inches(2.3),
        "当前实验配置",
        [
            "规则基线：full_open / random / periodic / round_robin / info_priority。",
            "学习型基线：dqn / cmdp_dqn / ppo。",
            "主任务目标：air_temperature_c、snow_surface_temperature_c、wind_speed_ms。",
            "评价指标：RMSE、MAE、sMAPE、Pearson、DTW。",
        ],
        fill=(240, 253, 244), line=(134, 239, 172),
    )
    add_bullet_box(
        slide,
        Inches(8.0), Inches(5.9), Inches(4.7), Inches(1.0),
        "一句话概括",
        ["这不是单独的预测任务，而是一条“调度 -> 状态估计 -> 预测 -> aggregate 评估”的完整链路。"],
        fill=(255, 247, 237), line=(253, 186, 116),
    )
    add_caption(slide, Inches(0.6), Inches(6.95), Inches(7.0), "图：当前 windblown 实验的完整工程架构。")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "结果总览：规则基线仍然最强，PPO 已进入可用区间")
    tradeoff = POSTHOC / "power_saving_vs_rmse_increase.png"
    task_focus = POSTHOC / "task_focus_primary" / "primary_target_rmse_increase_h1.png"
    slide.shapes.add_picture(str(tradeoff), Inches(0.55), Inches(1.0), width=Inches(6.25))
    slide.shapes.add_picture(str(task_focus), Inches(6.95), Inches(1.0), width=Inches(5.75))

    top = summary.set_index("scheduler")
    primary_top = primary.set_index("scheduler")
    bullets = [
        f"总体 aggregate：round_robin RMSE {top.loc['round_robin','rmse_increase_pct_vs_full_open']:+.2f}%，info_priority {top.loc['info_priority','rmse_increase_pct_vs_full_open']:+.2f}%，periodic {top.loc['periodic','rmse_increase_pct_vs_full_open']:+.2f}%。",
        f"PPO 当前总体 RMSE {top.loc['ppo','rmse_increase_pct_vs_full_open']:+.2f}%，节电 {top.loc['ppo','power_saving_pct_vs_full_open']:.1f}%。",
        f"主任务口径下，PPO 为 RMSE {primary_top.loc['ppo','rmse_increase_pct_vs_full_open']:+.2f}%，DTW {primary_top.loc['ppo','dtw_increase_pct_vs_full_open']:+.2f}%。",
        "当前 strongest baselines 仍是 round_robin / info_priority / periodic，说明规则方法非常强。",
    ]
    add_bullet_box(
        slide,
        Inches(0.7), Inches(5.35), Inches(12.0), Inches(1.45),
        "这页图想说明什么",
        bullets,
        fill=(248, 250, 252), line=(203, 213, 225),
    )
    add_caption(slide, Inches(0.6), Inches(4.95), Inches(5.8), "左：节电与 RMSE trade-off。")
    add_caption(slide, Inches(6.95), Inches(4.95), Inches(5.5), "右：主任务 H=1 的 RMSE 变化。")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "图例解读：PPO 已经在真实调度，但仍未超过最强规则基线")
    timeline = POSTHOC / "sensor_timelines" / "ppo" / "air_temperature_c_0_300_activation.png"
    curve = POSTHOC / "prediction_curves" / "informer" / "air_temperature_c_H1_overlay.png"
    heatmap = POSTHOC / "scheduler_model_rmse_increase_vs_full_open.png"
    slide.shapes.add_picture(str(timeline), Inches(0.5), Inches(1.0), width=Inches(6.1))
    slide.shapes.add_picture(str(curve), Inches(6.85), Inches(1.0), width=Inches(5.95))
    slide.shapes.add_picture(str(heatmap), Inches(0.85), Inches(4.45), width=Inches(4.8))
    add_bullet_box(
        slide,
        Inches(5.95), Inches(4.45), Inches(6.7), Inches(1.95),
        "三类图分别告诉我们什么",
        [
            "时间线图：PPO 已不是近固定策略，而是在两个高功耗传感器之间动态切换。",
            "预测曲线图：PPO 的曲线形态已经合理，但与最优规则基线相比仍有差距。",
            "热力图：规则基线跨模型更稳；PPO 作为开源 RL baseline 已可用，但还不是当前最优。",
        ],
        fill=(254, 242, 242), line=(252, 165, 165),
    )
    add_bullet_box(
        slide,
        Inches(5.95), Inches(6.55), Inches(6.7), Inches(0.65),
        "这次总结的最终落点",
        ["框架已经能稳定工作；下一步应把 cmdp_dqn 作为真正的方法创新主线。"],
        fill=(237, 233, 254), line=(196, 181, 253),
    )
    add_caption(slide, Inches(0.85), Inches(6.95), Inches(4.8), "左下：scheduler × predictor 的 RMSE 热力图。")

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build_deck()
