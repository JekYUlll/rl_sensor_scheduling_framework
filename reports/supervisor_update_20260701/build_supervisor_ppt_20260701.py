from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUT_DIR = Path(__file__).resolve().parent
FIG_DIR = OUT_DIR / "figures"
ASSET_DIR = OUT_DIR / "ppt_assets"
PPTX = OUT_DIR / "supervisor_update_easy_20260701.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COL = {
    "bg": "F7F3EA", "ink": "192334", "muted": "566273", "deep": "203A43",
    "green": "27735E", "green2": "DDEFE7", "orange": "B85F35", "orange2": "F5E1D2",
    "blue": "346E91", "blue2": "E2EEF4", "card": "FFFDF8", "line": "D8CEC1",
    "white": "FFFFFF", "gray": "EEE8DF", "red": "9B3A2D"
}

def rgb(h):
    h = h.strip('#')
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))

def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COL['bg'])

def add_text(slide, text, x, y, w, h, size: float = 13.0, bold=False, color='ink', align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.05
        run = p.add_run(); run.text = line
        run.font.name = 'Noto Sans CJK SC'
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(COL[color])
    return tb

def card(slide, x, y, w, h, fill='card', line='line', radius=True):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(COL[fill])
    s.line.color.rgb = rgb(COL[line]); s.line.width = Pt(0.8)
    return s

def pill(slide, text, x, y, w, fill, size=10.5):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.32))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(COL[fill]); s.line.fill.background()
    tf=s.text_frame; tf.clear(); tf.margin_top=Inches(0.03); tf.margin_bottom=Inches(0.02)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=True; r.font.color.rgb=rgb(COL['white'])
    return s

def picture(slide, path, x, y, w, h):
    from PIL import Image
    im=Image.open(path); iw,ih=im.size
    br=w/h; ir=iw/ih
    if ir>br:
        nw=w; nh=w/ir; px=x; py=y+(h-nh)/2
    else:
        nh=h; nw=h*ir; px=x+(w-nw)/2; py=y
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(nw), Inches(nh))

def footer(slide):
    add_text(slide, 'PD-PPO 实验更新 · 2026-07-01', 0.5, 7.12, 4.8, 0.18, size=8.5, color='muted')

# Slide 1
s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s, '实验主线从“传感器组合”改成“有限电量下选专家”', 0.45, 0.30, 11.8, 0.45, size=25, bold=True, color='deep')
add_text(s, '核心变化：把固定方案的解释空间压小，让“动态调度是否真的有用”变得可检验。', 0.48, 0.82, 12.1, 0.32, size=12.8, color='muted')

cols=[0.45,4.55,8.65]
fills=['orange2','green2','blue2']; accents=['orange','green','blue']
titles=['原来的实验容易被这样解释', '新实验更接近真实低功耗站点', '验证方式也更严格']
texts=[
    '多个传感器一起竞争（有些版本里，固定开某一组已经很强）。\n导师可能会问：RL 是在“按天气调度”，还是只找到了一个固定好组合？',
    '基础气象一直保留（低功耗背景）；高功耗专家仪器一次只能启用一台（激光、FC4、表面红外三选一）。\n不同天气事件需要不同专家，固定组合不再容易解释全部结果。',
    '同一批 24 次重复实验比较：固定方案、人工规则、PD-PPO。\n再看策略行为（是否随天气切换），排除“一直开一个”或“按顺序轮换”。'
]
for x,fill,acc,t,t2 in zip(cols,fills,accents,titles,texts):
    card(s,x,1.38,3.78,2.42,fill=fill)
    pill(s,t,x+0.18,1.56,2.25,acc,size=10)
    add_text(s,t2,x+0.22,2.02,3.34,1.35,size=11.5,color='ink')
for x in [4.20,8.30]:
    add_text(s,'→',x,2.23,0.25,0.35,size=22,bold=True,color='muted',align=PP_ALIGN.CENTER)

card(s,0.45,4.10,12.15,1.20,fill='card')
add_text(s,'通俗类比',0.70,4.27,1.1,0.25,size=12,bold=True,color='green')
add_text(s,'旧实验像“桌上摆多台仪器，看哪个组合好”（固定组合可能已经足够好）；新实验像“野外站电量有限”，基础气象常开，专家仪器只能按天气启用一台。',1.78,4.25,10.2,0.48,size=11.8,color='ink')
add_text(s,'因此这轮结果回答的是：模型能不能根据天气状态，决定今天该启用激光、FC4，还是表面红外。',1.78,4.82,10.2,0.26,size=11.8,bold=True,color='deep')

metrics=[('24/24','重复实验都保持正向优势\n（不是某一个随机种子好看）'),('+0.1494','逐时预测误差优势\n（正数表示 PD-PPO 更准）'),('+0.0710','按事件分开算后仍有优势\n（避免常见天气掩盖少见事件）')]
for i,(num,lab) in enumerate(metrics):
    x=0.45+i*4.08
    card(s,x,5.62,3.78,0.96,fill='white')
    add_text(s,num,x+0.22,5.77,1.35,0.35,size=22,bold=True,color='green')
    add_text(s,lab,x+1.48,5.73,2.04,0.43,size=10.3,color='ink')
footer(s)

# Slide 2
s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s,'新结果：不是单次运气，而是 24 次重复都站得住',0.45,0.30,11.4,0.45,size=25,bold=True,color='deep')
add_text(s,'本页聚焦稳定性：每个点对应一轮独立重复实验；数值为正表示 PD-PPO 预测更准。',0.48,0.82,12.2,0.32,size=12.5,color='muted')
card(s,0.45,1.20,12.05,3.48,fill='white')
picture(s,ASSET_DIR/'main_24seed_left_panel.png',0.82,1.35,7.85,3.05)
add_text(s,'图的要点：24 次重复的主要优势均为正；不是少数异常 seed 把平均值拉高。',0.78,4.31,7.9,0.22,size=10.8,color='muted')
card(s,8.92,1.46,2.95,0.72,fill='green2')
add_text(s,'24 个点均为正',9.10,1.62,2.55,0.24,size=14,bold=True,color='green',align=PP_ALIGN.CENTER)
add_text(s,'每个点是一轮独立重复',9.10,1.90,2.55,0.16,size=9.5,color='muted',align=PP_ALIGN.CENTER)
card(s,8.92,2.42,2.95,0.72,fill='orange2')
add_text(s,'最小优势仍为正',9.10,2.58,2.55,0.24,size=14,bold=True,color='orange',align=PP_ALIGN.CENTER)
add_text(s,'避免只报平均值好看',9.10,2.86,2.55,0.16,size=9.5,color='muted',align=PP_ALIGN.CENTER)
card(s,8.92,3.38,2.95,0.72,fill='blue2')
add_text(s,'固定对照被锁死',9.10,3.54,2.55,0.24,size=14,bold=True,color='blue',align=PP_ALIGN.CENTER)
add_text(s,'测试时不允许隐式轮换',9.10,3.82,2.55,0.16,size=9.5,color='muted',align=PP_ALIGN.CENTER)

# bottom summary cards
bottom=[
    ('稳定性结论', 'PD-PPO 全部 24 次重复实验都优于主要对照（固定方案、人工规则）。'),
    ('排除隐式切换', '固定方案全程同一组传感器，不能在测试时变成轮换策略。'),
    ('关键数值', '平均优势 0.1494；相对真正固定方案平均优势 0.0769（最小仍为正）。')
]
for i,(h,b) in enumerate(bottom):
    x=0.45+i*4.05
    card(s,x,5.00,3.78,1.05,fill=['green2','orange2','blue2'][i])
    add_text(s,h,x+0.18,5.15,1.5,0.23,size=11.6,bold=True,color=['green','orange','blue'][i])
    add_text(s,b,x+0.18,5.46,3.34,0.34,size=10.6,color='ink')
card(s,0.45,6.25,12.05,0.48,fill='card')
add_text(s,'推荐口径：不是证明“普通 PPO 普遍最强”；更准确地说，是证明在低功耗、专家仪器互斥的预测任务里，强化学习能稳定把有限电量分给更有用的仪器。',0.68,6.33,11.25,0.24,size=10.6,color='ink')
footer(s)

# Slide 3
s=prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_text(s,'行为证据：模型确实在按天气换专家',0.45,0.30,10.6,0.45,size=25,bold=True,color='deep')
add_text(s,'比单纯报误差更关键的是：它有没有学到符合物理直觉的开关方式。',0.48,0.82,12.2,0.32,size=12.5,color='muted')
card(s,0.45,1.20,7.02,3.78,fill='white')
picture(s,ASSET_DIR/'behavior_heatmap_panel.png',0.72,1.38,6.20,3.22)
add_text(s,'左图看选择频率：颜色越深，表示该天气下越常开这台仪器。',0.72,4.55,6.4,0.22,size=10.4,color='muted')

# concise right explanation
blocks=[
    ('不是固定启用同一台', '如果只是固定方案，三类天气下会反复选同一台仪器。现在不是这样。'),
    ('选择模式符合物理直觉', '粒子→激光：看颗粒；通量→FC4：看吹雪通量；热状态→表面红外：看地表热。'),
    ('消融检查', '去掉事件状态信息后，严格固定对照从 24/24 降到 21/24（稳定性下降）。')
]
for i,(h,b) in enumerate(blocks):
    y=1.20+i*1.20
    card(s,7.70,y,4.82,0.96,fill=['orange2','green2','blue2'][i])
    add_text(s,h,7.92,y+0.13,1.45,0.22,size=11.5,bold=True,color=['orange','green','blue'][i])
    add_text(s,b,9.20,y+0.12,3.02,0.42,size=10.4,color='ink')

card(s,0.45,5.33,12.07,0.92,fill='card')
add_text(s,'可直接汇报的一句话',0.68,5.48,1.75,0.22,size=11.8,bold=True,color='deep')
add_text(s,'当前实验更接近真实低功耗站：基础气象常开，高功耗专家仪器按天气事件轮流启用。结果显示，PD-PPO 不只是降低误差，还学到和物理事件对应的开关策略。',2.38,5.44,9.72,0.30,size=10.8,color='ink')
add_text(s,'适用边界：结论限于当前微气候仿真和这一类“专家仪器预算调度”问题；不外推为所有传感器调度都普遍最优。',0.68,6.48,11.6,0.24,size=10.4,color='muted')
footer(s)

prs.core_properties.title='PD-PPO 实验更新导师简报'
prs.core_properties.subject='Plain-language supervisor update for 2026-07-01 PD-PPO experiment mainline'
prs.core_properties.author='Hermes'
prs.save(PPTX)
print(PPTX)
